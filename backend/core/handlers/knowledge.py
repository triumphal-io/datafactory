import chromadb
from chromadb.utils import embedding_functions
import pandas as pd
import os
from pathlib import Path
from django.conf import settings

# ChromaDB client configuration
CHROMA_DB_PATH = os.path.join(settings.BASE_DIR, 'storage', 'chromadb')
_chroma_client = None  # Lazy-loaded singleton


def get_chroma_client():
    """
    Get or create ChromaDB client (lazy initialization).
    Only creates the connection when first needed.
    """
    global _chroma_client
    if _chroma_client is None:
        Path(CHROMA_DB_PATH).mkdir(parents=True, exist_ok=True)
        _chroma_client = chromadb.PersistentClient(path=CHROMA_DB_PATH)
    return _chroma_client


def get_embedding_function():
    """
    Get the configured embedding function based on settings.
    
    Returns:
        EmbeddingFunction: Configured embedding function for ChromaDB
    """
    embedding_type = getattr(settings, 'EMBEDDING_MODEL_TYPE', 'default')
    
    if embedding_type == 'openai':
        model_name = getattr(settings, 'EMBEDDING_MODEL_NAME', 'text-embedding-3-small')
        return embedding_functions.OpenAIEmbeddingFunction(
            api_key=os.getenv("OPENAI_API_KEY"),
            model_name=model_name
        )
    elif embedding_type == 'sentence-transformers':
        model_name = getattr(settings, 'EMBEDDING_MODEL_NAME', 'all-MiniLM-L6-v2')
        return embedding_functions.SentenceTransformerEmbeddingFunction(
            model_name=model_name
        )
    elif embedding_type == 'default':
        # Use ChromaDB's default embedding function
        return embedding_functions.DefaultEmbeddingFunction()
    else:
        raise ValueError(f"Unsupported embedding model type: {embedding_type}")


def get_or_create_collection(collection_name="file_chunks"):
    """
    Get or create a ChromaDB collection for storing file chunks.
    Uses the embedding model configured in settings.
    
    Args:
        collection_name (str): Name of the collection
        
    Returns:
        Collection: ChromaDB collection instance
    """
    embedding_func = get_embedding_function()
    embedding_type = getattr(settings, 'EMBEDDING_MODEL_TYPE', 'default')
    embedding_model_name = getattr(settings, 'EMBEDDING_MODEL_NAME', 'default')
    
    try:
        collection = get_chroma_client().get_collection(
            name=collection_name,
            embedding_function=embedding_func
        )
    except ValueError as e:
        # If collection exists with different embedding function, delete and recreate
        if "already exists" in str(e).lower():
            print(f"⚠ Collection exists with different embedding function. Recreating with {embedding_type} embeddings...")
            try:
                get_chroma_client().delete_collection(name=collection_name)
            except:
                pass
            collection = get_chroma_client().create_collection(
                name=collection_name,
                embedding_function=embedding_func,
                metadata={
                    "description": "File chunks for RAG",
                    "embedding_type": embedding_type,
                    "embedding_model": embedding_model_name
                }
            )
        else:
            raise
    except Exception:
        # Collection doesn't exist, create it
        collection = get_chroma_client().create_collection(
            name=collection_name,
            embedding_function=embedding_func,
            metadata={
                "description": "File chunks for RAG",
                "embedding_type": embedding_type,
                "embedding_model": embedding_model_name
            }
        )
    return collection


def index_csv_file(file_path, file_id, filename, user_id, folder_name=None):
    """
    Index a CSV file by creating chunks for each row and storing in ChromaDB.
    
    Args:
        file_path (str): Path to the CSV file
        file_id (str): UUID of the file in database
        filename (str): Original filename
        user_id (int): ID of the user who uploaded the file
        folder_name (str, optional): Name of the folder containing the file
        
    Returns:
        int: Number of chunks created
    """
    try:
        # Read CSV file
        df = pd.read_csv(file_path)
        
        chunks = []
        metadatas = []
        ids = []
        
        # Batch size to avoid exceeding token limits (100 rows per batch)
        BATCH_SIZE = 100
        collection = get_or_create_collection()
        total_chunks = 0
        
        for idx, row in df.iterrows():
            # Create text chunk for this row
            chunk = f"Record {idx + 1}: "
            for col in df.columns:
                if pd.notna(row[col]):
                    chunk += f"{col} is {row[col]}. "
            
            chunks.append(chunk)
            
            # Create metadata for this chunk
            metadata = {
                'row_id': int(idx),
                'source': filename,
                'file_id': str(file_id),
                'user_id': str(user_id),
                'file_type': 'csv'
            }
            
            # Add folder name if provided
            if folder_name:
                metadata['folder_name'] = folder_name
            
            # Add first 3 columns as metadata for better filtering
            for col_idx, col in enumerate(df.columns[:3]):
                if pd.notna(row[col]):
                    metadata[col] = str(row[col])
            
            metadatas.append(metadata)
            ids.append(f"{file_id}_row_{idx}")
            
            # Add batch to ChromaDB when batch size is reached
            if len(chunks) >= BATCH_SIZE:
                collection.add(
                    documents=chunks,
                    metadatas=metadatas,
                    ids=ids
                )
                total_chunks += len(chunks)
                chunks = []
                metadatas = []
                ids = []
        
        # Add remaining chunks
        if chunks:
            collection.add(
                documents=chunks,
                metadatas=metadatas,
                ids=ids
            )
            total_chunks += len(chunks)
        
        print(f"✓ Indexed {total_chunks} chunks from CSV file: {filename}")
        return total_chunks
        
    except Exception as e:
        print(f"✗ Error indexing CSV file {filename}: {str(e)}")
        raise


def index_xlsx_file(file_path, file_id, filename, user_id, folder_name=None):
    """
    Index an XLSX file by creating chunks for each row across all sheets.
    
    Args:
        file_path (str): Path to the XLSX file
        file_id (str): UUID of the file in database
        filename (str): Original filename
        user_id (int): ID of the user who uploaded the file
        folder_name (str, optional): Name of the folder containing the file
        
    Returns:
        int: Number of chunks created
    """
    try:
        # Read all sheets from XLSX file
        xlsx_file = pd.ExcelFile(file_path)
        
        total_chunks = 0
        all_chunks = []
        all_metadatas = []
        all_ids = []
        
        # Batch size to avoid exceeding token limits (100 rows per batch)
        BATCH_SIZE = 100
        collection = get_or_create_collection()
        
        for sheet_name in xlsx_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            
            for idx, row in df.iterrows():
                # Create text chunk for this row
                chunk = f"Record {idx + 1}"
                if len(xlsx_file.sheet_names) > 1:
                    chunk += f" (Sheet: {sheet_name})"
                chunk += ": "
                
                for col in df.columns:
                    if pd.notna(row[col]):
                        chunk += f"{col} is {row[col]}. "
                
                all_chunks.append(chunk)
                
                # Create metadata for this chunk
                metadata = {
                    'row_id': int(idx),
                    'source': filename,
                    'sheet_name': sheet_name,
                    'file_id': str(file_id),
                    'user_id': str(user_id),
                    'file_type': 'xlsx'
                }
                
                # Add folder name if provided
                if folder_name:
                    metadata['folder_name'] = folder_name
                
                # Add first 3 columns as metadata for better filtering
                for col_idx, col in enumerate(df.columns[:3]):
                    if pd.notna(row[col]):
                        metadata[col] = str(row[col])
                
                all_metadatas.append(metadata)
                all_ids.append(f"{file_id}_sheet_{sheet_name}_row_{idx}")
                
                # Add batch to ChromaDB when batch size is reached
                if len(all_chunks) >= BATCH_SIZE:
                    collection.add(
                        documents=all_chunks,
                        metadatas=all_metadatas,
                        ids=all_ids
                    )
                    total_chunks += len(all_chunks)
                    all_chunks = []
                    all_metadatas = []
                    all_ids = []
        
        # Add remaining chunks
        if all_chunks:
            collection.add(
                documents=all_chunks,
                metadatas=all_metadatas,
                ids=all_ids
            )
            total_chunks += len(all_chunks)
        
        print(f"✓ Indexed {total_chunks} chunks from XLSX file: {filename}")
        return total_chunks
        
    except Exception as e:
        print(f"✗ Error indexing XLSX file {filename}: {str(e)}")
        raise


def index_pdf_file(file_path, file_id, filename, user_id, folder_name=None):
    """
    Index a PDF file by creating chunks for each page.
    
    Args:
        file_path (str): Path to the PDF file
        file_id (str): UUID of the file in database
        filename (str): Original filename
        user_id (int): ID of the user who uploaded the file
        folder_name (str, optional): Name of the folder containing the file
        
    Returns:
        int: Number of chunks created
    """
    try:
        import PyPDF2
        
        chunks = []
        metadatas = []
        ids = []
        
        collection = get_or_create_collection()
        total_chunks = 0
        
        with open(file_path, 'rb') as pdf_file:
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            num_pages = len(pdf_reader.pages)
            
            for page_num in range(num_pages):
                page = pdf_reader.pages[page_num]
                text = page.extract_text()
                
                if not text.strip():
                    continue
                
                # Create text chunk for this page
                chunk = f"Page {page_num + 1}: {text.strip()}"
                chunks.append(chunk)
                
                # Create metadata for this chunk
                metadata = {
                    'page_num': page_num + 1,
                    'source': filename,
                    'file_id': str(file_id),
                    'user_id': str(user_id),
                    'file_type': 'pdf'
                }
                
                # Add folder name if provided
                if folder_name:
                    metadata['folder_name'] = folder_name
                
                metadatas.append(metadata)
                ids.append(f"{file_id}_page_{page_num}")
        
        # Add all chunks to ChromaDB
        if chunks:
            collection.add(
                documents=chunks,
                metadatas=metadatas,
                ids=ids
            )
            total_chunks = len(chunks)
        
        print(f"✓ Indexed {total_chunks} chunks from PDF file: {filename}")
        return total_chunks
        
    except Exception as e:
        print(f"✗ Error indexing PDF file {filename}: {str(e)}")
        raise


def index_docx_file(file_path, file_id, filename, user_id, folder_name=None):
    """
    Index a DOCX file by creating chunks for paragraphs and tables.
    
    Args:
        file_path (str): Path to the DOCX file
        file_id (str): UUID of the file in database
        filename (str): Original filename
        user_id (int): ID of the user who uploaded the file
        folder_name (str, optional): Name of the folder containing the file
        
    Returns:
        int: Number of chunks created
    """
    try:
        from docx import Document as DocxDocument
        
        chunks = []
        metadatas = []
        ids = []
        
        collection = get_or_create_collection()
        total_chunks = 0
        chunk_num = 0
        
        doc = DocxDocument(file_path)
        
        # Index paragraphs
        current_section = []
        for para_idx, paragraph in enumerate(doc.paragraphs):
            text = paragraph.text.strip()
            if not text:
                continue
            
            # If this is a heading, save previous section and start new one
            if paragraph.style and paragraph.style.name and paragraph.style.name.startswith('Heading'):
                # Save previous section if it exists
                if current_section:
                    chunk = "\n".join(current_section)
                    chunks.append(chunk)
                    
                    metadata = {
                        'chunk_num': chunk_num,
                        'source': filename,
                        'file_id': str(file_id),
                        'user_id': str(user_id),
                        'file_type': 'docx',
                        'content_type': 'paragraph'
                    }
                    
                    if folder_name:
                        metadata['folder_name'] = folder_name
                    
                    metadatas.append(metadata)
                    ids.append(f"{file_id}_chunk_{chunk_num}")
                    chunk_num += 1
                    current_section = []
                
                # Start new section with heading
                current_section.append(text)
            else:
                current_section.append(text)
            
            # Create chunks of reasonable size (every 5 paragraphs)
            if len(current_section) >= 5:
                chunk = "\n".join(current_section)
                chunks.append(chunk)
                
                metadata = {
                    'chunk_num': chunk_num,
                    'source': filename,
                    'file_id': str(file_id),
                    'user_id': str(user_id),
                    'file_type': 'docx',
                    'content_type': 'paragraph'
                }
                
                if folder_name:
                    metadata['folder_name'] = folder_name
                
                metadatas.append(metadata)
                ids.append(f"{file_id}_chunk_{chunk_num}")
                chunk_num += 1
                current_section = []
        
        # Save any remaining section
        if current_section:
            chunk = "\n".join(current_section)
            chunks.append(chunk)
            
            metadata = {
                'chunk_num': chunk_num,
                'source': filename,
                'file_id': str(file_id),
                'user_id': str(user_id),
                'file_type': 'docx',
                'content_type': 'paragraph'
            }
            
            if folder_name:
                metadata['folder_name'] = folder_name
            
            metadatas.append(metadata)
            ids.append(f"{file_id}_chunk_{chunk_num}")
            chunk_num += 1
        
        # Index tables separately
        for table_idx, table in enumerate(doc.tables):
            table_text = f"Table {table_idx + 1}:\n"
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                table_text += " | ".join(cells) + "\n"
            
            chunks.append(table_text)
            
            metadata = {
                'table_num': table_idx + 1,
                'chunk_num': chunk_num,
                'source': filename,
                'file_id': str(file_id),
                'user_id': str(user_id),
                'file_type': 'docx',
                'content_type': 'table'
            }
            
            if folder_name:
                metadata['folder_name'] = folder_name
            
            metadatas.append(metadata)
            ids.append(f"{file_id}_table_{table_idx}")
            chunk_num += 1
        
        # Add all chunks to ChromaDB
        if chunks:
            collection.add(
                documents=chunks,
                metadatas=metadatas,
                ids=ids
            )
            total_chunks = len(chunks)
        
        print(f"✓ Indexed {total_chunks} chunks from DOCX file: {filename}")
        return total_chunks
        
    except Exception as e:
        print(f"✗ Error indexing DOCX file {filename}: {str(e)}")
        raise


def index_txt_file(file_path, file_id, filename, user_id, folder_name=None):
    """
    Index a TXT file by creating chunks for paragraphs (split by double newlines).
    
    Args:
        file_path (str): Path to the TXT file
        file_id (str): UUID of the file in database
        filename (str): Original filename
        user_id (int): ID of the user who uploaded the file
        folder_name (str, optional): Name of the folder containing the file
        
    Returns:
        int: Number of chunks created
    """
    try:
        chunks = []
        metadatas = []
        ids = []
        
        collection = get_or_create_collection()
        total_chunks = 0
        
        with open(file_path, 'r', encoding='utf-8') as txt_file:
            content = txt_file.read()
            
            # Split by paragraphs (double newlines or single newlines for simpler files)
            paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
            
            # If no double newlines found, split by single newlines and group every 5 lines
            if len(paragraphs) <= 1:
                lines = [line.strip() for line in content.split('\n') if line.strip()]
                paragraphs = []
                for i in range(0, len(lines), 5):
                    paragraphs.append('\n'.join(lines[i:i+5]))
            
            for chunk_num, paragraph in enumerate(paragraphs):
                if not paragraph:
                    continue
                
                chunks.append(paragraph)
                
                metadata = {
                    'chunk_num': chunk_num,
                    'source': filename,
                    'file_id': str(file_id),
                    'user_id': str(user_id),
                    'file_type': 'txt'
                }
                
                if folder_name:
                    metadata['folder_name'] = folder_name
                
                metadatas.append(metadata)
                ids.append(f"{file_id}_chunk_{chunk_num}")
        
        # Add all chunks to ChromaDB
        if chunks:
            collection.add(
                documents=chunks,
                metadatas=metadatas,
                ids=ids
            )
            total_chunks = len(chunks)
        
        print(f"✓ Indexed {total_chunks} chunks from TXT file: {filename}")
        return total_chunks
        
    except Exception as e:
        print(f"✗ Error indexing TXT file {filename}: {str(e)}")
        raise


def index_md_file(file_path, file_id, filename, user_id, folder_name=None):
    """
    Index a Markdown file by creating chunks for sections (split by headings).
    
    Args:
        file_path (str): Path to the MD file
        file_id (str): UUID of the file in database
        filename (str): Original filename
        user_id (int): ID of the user who uploaded the file
        folder_name (str, optional): Name of the folder containing the file
        
    Returns:
        int: Number of chunks created
    """
    try:
        chunks = []
        metadatas = []
        ids = []
        
        collection = get_or_create_collection()
        total_chunks = 0
        
        with open(file_path, 'r', encoding='utf-8') as md_file:
            content = md_file.read()
            lines = content.split('\n')
            
            current_section = []
            current_heading = None
            chunk_num = 0
            
            for line in lines:
                # Check if line is a heading
                if line.strip().startswith('#'):
                    # Save previous section if it exists
                    if current_section:
                        chunk_text = '\n'.join(current_section)
                        chunks.append(chunk_text)
                        
                        metadata = {
                            'chunk_num': chunk_num,
                            'source': filename,
                            'file_id': str(file_id),
                            'user_id': str(user_id),
                            'file_type': 'md'
                        }
                        
                        if current_heading:
                            metadata['heading'] = current_heading
                        if folder_name:
                            metadata['folder_name'] = folder_name
                        
                        metadatas.append(metadata)
                        ids.append(f"{file_id}_chunk_{chunk_num}")
                        chunk_num += 1
                        current_section = []
                    
                    # Start new section with heading
                    current_heading = line.strip().lstrip('#').strip()
                    current_section.append(line)
                else:
                    current_section.append(line)
            
            # Save last section
            if current_section:
                chunk_text = '\n'.join(current_section)
                chunks.append(chunk_text)
                
                metadata = {
                    'chunk_num': chunk_num,
                    'source': filename,
                    'file_id': str(file_id),
                    'user_id': str(user_id),
                    'file_type': 'md'
                }
                
                if current_heading:
                    metadata['heading'] = current_heading
                if folder_name:
                    metadata['folder_name'] = folder_name
                
                metadatas.append(metadata)
                ids.append(f"{file_id}_chunk_{chunk_num}")
        
        # Add all chunks to ChromaDB
        if chunks:
            collection.add(
                documents=chunks,
                metadatas=metadatas,
                ids=ids
            )
            total_chunks = len(chunks)
        
        print(f"✓ Indexed {total_chunks} chunks from MD file: {filename}")
        return total_chunks
        
    except Exception as e:
        print(f"✗ Error indexing MD file {filename}: {str(e)}")
        raise


def index_pptx_file(file_path, file_id, filename, user_id, folder_name=None):
    """
    Index a PPTX file by creating chunks for each slide.
    
    Args:
        file_path (str): Path to the PPTX file
        file_id (str): UUID of the file in database
        filename (str): Original filename
        user_id (int): ID of the user who uploaded the file
        folder_name (str, optional): Name of the folder containing the file
        
    Returns:
        int: Number of chunks created
    """
    try:
        from pptx import Presentation
        
        chunks = []
        metadatas = []
        ids = []
        
        collection = get_or_create_collection()
        total_chunks = 0
        
        prs = Presentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts = []
            slide_title = None
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    
                    # Capture title
                    if hasattr(shape, "shape_type") and shape == slide.shapes.title:
                        slide_title = text
                    
                    slide_texts.append(text)
                
                # Extract text from tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        slide_texts.append(" | ".join(cells))
            
            # Extract notes if present
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_texts.append(f"Notes: {notes_text}")
            
            # Create chunk if slide has content
            if slide_texts:
                chunk = f"Slide {slide_num}: " + " ".join(slide_texts)
                chunks.append(chunk)
                
                metadata = {
                    'slide_num': slide_num,
                    'source': filename,
                    'file_id': str(file_id),
                    'user_id': str(user_id),
                    'file_type': 'pptx'
                }
                
                if slide_title:
                    metadata['slide_title'] = slide_title
                if folder_name:
                    metadata['folder_name'] = folder_name
                
                metadatas.append(metadata)
                ids.append(f"{file_id}_slide_{slide_num}")
        
        # Add all chunks to ChromaDB
        if chunks:
            collection.add(
                documents=chunks,
                metadatas=metadatas,
                ids=ids
            )
            total_chunks = len(chunks)
        
        print(f"✓ Indexed {total_chunks} chunks from PPTX file: {filename}")
        return total_chunks
        
    except Exception as e:
        print(f"✗ Error indexing PPTX file {filename}: {str(e)}")
        raise


def index_file(file_path, file_id, filename, user_id, folder_name=None):
    """
    Index a file in ChromaDB based on its type.
    
    Args:
        file_path (str): Path to the file
        file_id (str): UUID of the file in database
        filename (str): Original filename
        user_id (int): ID of the user who uploaded the file
        folder_name (str, optional): Name of the folder containing the file
        
    Returns:
        int: Number of chunks created
    """
    file_extension = os.path.splitext(filename)[1].lower()
    
    if file_extension == '.csv':
        return index_csv_file(file_path, file_id, filename, user_id, folder_name)
    elif file_extension in ['.xlsx', '.xls']:
        return index_xlsx_file(file_path, file_id, filename, user_id, folder_name)
    elif file_extension == '.pdf':
        return index_pdf_file(file_path, file_id, filename, user_id, folder_name)
    elif file_extension in ['.docx', '.doc']:
        return index_docx_file(file_path, file_id, filename, user_id, folder_name)
    elif file_extension == '.txt':
        return index_txt_file(file_path, file_id, filename, user_id, folder_name)
    elif file_extension == '.md':
        return index_md_file(file_path, file_id, filename, user_id, folder_name)
    elif file_extension in ['.pptx', '.ppt']:
        return index_pptx_file(file_path, file_id, filename, user_id, folder_name)
    else:
        raise ValueError(f"Unsupported file type for indexing: {file_extension}")


def query_rag(query, filename=None, user_id=None, n_results=5, search_type='query'):
    """
    Query the RAG system to retrieve relevant chunks.
    
    Args:
        query (str): The search query
        filename (str, optional): Filter by specific filename
        user_id (int, optional): Filter by user ID to isolate user data
        n_results (int): Number of results to return (default: 5, ignored for identifier search)
        search_type (str): Type of search - 'identifier' for exact text matching (IDs, codes) - returns only top 1,
                          'query' for semantic search (default: 'query') - returns n_results
        
    Returns:
        dict: Query results with documents, metadatas, and distances (distances only for 'query' type)
    """
    try:
        collection = get_or_create_collection()
        
        # Build where clause for filtering with proper ChromaDB syntax
        where_clause = None
        conditions = []
        
        if user_id is not None:
            conditions.append({'user_id': str(user_id)})
        if filename is not None:
            conditions.append({'source': filename})
        
        # ChromaDB requires $and operator for multiple conditions
        if len(conditions) == 1:
            where_clause = conditions[0]
        elif len(conditions) > 1:
            where_clause = {'$and': conditions}
        
        # Use different search methods based on search_type
        if search_type == 'identifier':
            # Use ChromaDB's native text search for exact matching
            # Build combined where clause that includes document text matching
            where_document_clause = {"$contains": query}
            
            # Use get() for text-based filtering
            results = collection.get(
                where=where_clause,
                where_document=where_document_clause,
                limit=1  # Only return top 1 result for identifiers
            )
            
            # Format results to match expected output structure
            # Note: get() doesn't return distances, only query() does
            return {
                'documents': results['documents'] if results['documents'] else [],
                'metadatas': results['metadatas'] if results['metadatas'] else [],
                'distances': []  # No distances for text-based search
            }
        else:
            # Use semantic search for 'query' type
            results = collection.query(
                query_texts=[query],
                n_results=n_results,
                where=where_clause
            )
            
            # Extract and format results
            return {
                'documents': results['documents'][0] if results['documents'] else [],
                'metadatas': results['metadatas'][0] if results['metadatas'] else [],
                'distances': results['distances'][0] if results['distances'] else []
            }
        
    except Exception as e:
        print(f"✗ Error querying RAG: {str(e)}")
        return {
            'documents': [],
            'metadatas': [],
            'distances': [],
            'error': str(e)
        }


def delete_file_chunks(file_id):
    """
    Delete all chunks associated with a specific file.
    
    Args:
        file_id (str): UUID of the file to delete chunks for
    """
    try:
        collection = get_or_create_collection()
        
        # Delete all chunks with this file_id
        collection.delete(
            where={"file_id": str(file_id)}
        )
        
        print(f"✓ Deleted chunks for file: {file_id}")
        
    except Exception as e:
        print(f"✗ Error deleting file chunks: {str(e)}")


def get_user_files(user_id):
    """
    Get list of unique filenames for a specific user.
    
    Args:
        user_id (int): User ID
        
    Returns:
        list: List of unique filenames
    """
    try:
        collection = get_or_create_collection()
        
        # Get all chunks for this user
        results = collection.get(
            where={"user_id": str(user_id)},
            include=["metadatas"]
        )
        
        # Extract unique filenames
        filenames = set()
        if results and results.get('metadatas'):
            for metadata in results['metadatas']:
                if 'source' in metadata:
                    filenames.add(metadata['source'])
        
        return sorted(list(filenames))
        
    except Exception as e:
        print(f"✗ Error getting user files: {str(e)}")
        return []


def update_file_metadata(file_id, new_filename):
    """
    Update the filename in metadata for all chunks of a file.
    
    Args:
        file_id (str): UUID of the file
        new_filename (str): New filename to set
    """
    try:
        collection = get_or_create_collection()
        
        # Get all chunks for this file
        results = collection.get(
            where={"file_id": str(file_id)},
            include=["metadatas"]
        )
        
        if results and results.get('ids'):
            # Update each chunk's metadata
            for idx, chunk_id in enumerate(results['ids']):
                metadata = results['metadatas'][idx]
                metadata['source'] = new_filename
                
                # Update the chunk with new metadata
                collection.update(
                    ids=[chunk_id],
                    metadatas=[metadata]
                )
        
        print(f"✓ Updated filename metadata for file: {file_id}")
        
    except Exception as e:
        print(f"✗ Error updating file metadata: {str(e)}")
        raise


def update_folder_metadata(old_folder_name, new_folder_name, user_id):
    """
    Update the folder name in metadata for all chunks in a folder.
    
    Args:
        old_folder_name (str): Current folder name
        new_folder_name (str): New folder name
        user_id (int): User ID to scope the update
    """
    try:
        collection = get_or_create_collection()
        
        # Get all chunks for this folder
        results = collection.get(
            where={
                "$and": [
                    {"user_id": str(user_id)},
                    {"folder_name": old_folder_name}
                ]
            },
            include=["metadatas"]
        )
        
        if results and results.get('ids'):
            # Update each chunk's metadata
            for idx, chunk_id in enumerate(results['ids']):
                metadata = results['metadatas'][idx]
                metadata['folder_name'] = new_folder_name
                
                # Update the chunk with new metadata
                collection.update(
                    ids=[chunk_id],
                    metadatas=[metadata]
                )
        
        print(f"✓ Updated folder metadata from '{old_folder_name}' to '{new_folder_name}'")
        
    except Exception as e:
        print(f"✗ Error updating folder metadata: {str(e)}")
        raise


def delete_folder_chunks(folder_name, user_id):
    """
    Delete all chunks associated with files in a specific folder.
    
    Args:
        folder_name (str): Name of the folder
        user_id (int): User ID to scope the deletion
    """
    try:
        collection = get_or_create_collection()
        
        # Delete all chunks with this folder_name and user_id
        collection.delete(
            where={
                "$and": [
                    {"user_id": str(user_id)},
                    {"folder_name": folder_name}
                ]
            }
        )
        
        print(f"✓ Deleted chunks for folder: {folder_name}")
        
    except Exception as e:
        print(f"✗ Error deleting folder chunks: {str(e)}")
