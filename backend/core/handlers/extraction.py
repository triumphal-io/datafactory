import csv
import os
import threading
from django.conf import settings
from core.models import File
from openpyxl import load_workbook
import PyPDF2
from docx import Document as DocxDocument


class FileExtractor:
    """
    A class to handle extraction and conversion of various file types to markdown format.
    Supports CSV, XLSX, and can be extended for PDF, PPT, images, etc.
    """
    
    def __init__(self, file_path, filename):
        """
        Initialize the FileExtractor.
        
        Args:
            file_path: Path to the file to extract
            filename: Original filename (used to determine file type)
        """
        self.file_path = file_path
        self.filename = filename
        self.file_extension = os.path.splitext(filename)[1].lower()
    
    def extract(self):
        """
        Main extraction method that routes to appropriate converter based on file type.
        
        Returns:
            Markdown formatted string representation of the file content
        """
        try:
            if self.file_extension == '.csv':
                return self.csv_to_md()
            elif self.file_extension in ['.xlsx', '.xls']:
                return self.xlsx_to_md()
            elif self.file_extension == '.pdf':
                return self.pdf_to_md()
            elif self.file_extension in ['.docx', '.doc']:
                return self.docx_to_md()
            elif self.file_extension == '.txt':
                return self.txt_to_md()
            elif self.file_extension == '.md':
                return self.md_to_md()
            elif self.file_extension in ['.pptx', '.ppt']:
                return self.pptx_to_md()
            # Future file types can be added here:
            # elif self.file_extension in ['.png', '.jpg', '.jpeg']:
            #     return self.image_to_md()
            else:
                return f"Unsupported file type: {self.file_extension}"
        except Exception as e:
            return f"Error extracting {self.file_extension}: {str(e)}"
    
    def csv_to_md(self):
        """
        Convert CSV file to markdown table format.
        
        Returns:
            Markdown formatted string representation of the CSV
        """
        try:
            with open(self.file_path, 'r', encoding='utf-8') as csvfile:
                reader = csv.reader(csvfile)
                rows = list(reader)
                
                if not rows:
                    return "Empty CSV file"
                
                # Build markdown table
                markdown_lines = []
                
                # Header row
                if len(rows) > 0:
                    header = rows[0]
                    # Replace newlines with HTML line breaks to prevent table breaking
                    header = [cell.replace('\r\n', '<br>').replace('\n', '<br>').replace('\r', '<br>') for cell in header]
                    markdown_lines.append('| ' + ' | '.join(header) + ' |')
                    markdown_lines.append('| ' + ' | '.join(['---'] * len(header)) + ' |')
                
                # Data rows
                for row in rows[1:]:
                    # Pad row if it has fewer columns than header
                    if len(row) < len(header):
                        row = row + [''] * (len(header) - len(row))
                    # Truncate if it has more columns
                    elif len(row) > len(header):
                        row = row[:len(header)]
                    # Replace newlines with HTML line breaks to prevent table breaking
                    row = [cell.replace('\r\n', '<br>').replace('\n', '<br>').replace('\r', '<br>') for cell in row]
                    markdown_lines.append('| ' + ' | '.join(row) + ' |')
                
                return '\n'.join(markdown_lines)
                
        except Exception as e:
            raise Exception(f"CSV extraction failed: {str(e)}")
    
    def xlsx_to_md(self):
        """
        Convert XLSX file to markdown table format.
        Calculates all formulas and displays them as text.
        If there are multiple sheets, they are combined with sheet names as headings.
        
        Returns:
            Markdown formatted string representation of the XLSX
        """
        try:
            workbook = load_workbook(self.file_path, data_only=True)
            markdown_sections = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                # Add sheet name as heading for multiple sheets
                if len(workbook.sheetnames) > 1:
                    markdown_sections.append(f"## {sheet_name}\n")
                
                # Get all rows with data
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    # Convert all cells to strings, handling None values and newlines
                    row_data = [str(cell).replace('\r\n', '<br>').replace('\n', '<br>').replace('\r', '<br>') if cell is not None else '' for cell in row]
                    # Skip completely empty rows
                    if any(cell for cell in row_data):
                        rows.append(row_data)
                
                if not rows:
                    markdown_sections.append("*Empty sheet*\n")
                    continue
                
                # Determine the maximum number of columns
                max_cols = max(len(row) for row in rows) if rows else 0
                
                # Normalize all rows to have the same number of columns
                for i in range(len(rows)):
                    if len(rows[i]) < max_cols:
                        rows[i] = rows[i] + [''] * (max_cols - len(rows[i]))
                
                # Build markdown table
                markdown_lines = []
                
                if rows:
                    # First row as header
                    header = rows[0]
                    # Escape pipe characters in cells
                    header = [cell.replace('|', '\\|') for cell in header]
                    markdown_lines.append('| ' + ' | '.join(header) + ' |')
                    markdown_lines.append('| ' + ' | '.join(['---'] * len(header)) + ' |')
                    
                    # Data rows
                    for row in rows[1:]:
                        # Escape pipe characters in cells
                        row = [cell.replace('|', '\\|') for cell in row]
                        markdown_lines.append('| ' + ' | '.join(row) + ' |')
                
                markdown_sections.append('\n'.join(markdown_lines))
            
            # Join all sections with double newlines
            return '\n\n'.join(markdown_sections)
            
        except Exception as e:
            raise Exception(f"XLSX extraction failed: {str(e)}")
    
    def pdf_to_md(self):
        """
        Convert PDF file to markdown format.
        Extracts text content from all pages.
        
        Returns:
            Markdown formatted string representation of the PDF
        """
        try:
            markdown_sections = []
            
            with open(self.file_path, 'rb') as pdf_file:
                pdf_reader = PyPDF2.PdfReader(pdf_file)
                num_pages = len(pdf_reader.pages)
                
                for page_num in range(num_pages):
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()
                    
                    if text.strip():
                        # Add page header for multi-page PDFs
                        if num_pages > 1:
                            markdown_sections.append(f"## Page {page_num + 1}\n")
                        markdown_sections.append(text.strip())
                
                if not markdown_sections:
                    return "*Empty or unreadable PDF*"
                
                return "\n\n".join(markdown_sections)
                
        except Exception as e:
            raise Exception(f"PDF extraction failed: {str(e)}")
    
    def docx_to_md(self):
        """
        Convert DOCX file to markdown format.
        Preserves headings, paragraphs, and basic formatting.
        
        Returns:
            Markdown formatted string representation of the DOCX
        """
        try:
            doc = DocxDocument(self.file_path)
            markdown_lines = []
            
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue
                
                # Convert heading styles to markdown headings
                if paragraph.style and paragraph.style.name and paragraph.style.name.startswith('Heading'):
                    try:
                        level = int(paragraph.style.name.replace('Heading ', ''))
                        markdown_lines.append(f"{'#' * level} {text}")
                    except (ValueError, AttributeError):
                        markdown_lines.append(text)
                else:
                    markdown_lines.append(text)
                
                markdown_lines.append("")  # Add blank line between paragraphs
            
            # Handle tables
            if doc.tables:
                for table_num, table in enumerate(doc.tables):
                    if markdown_lines and markdown_lines[-1] != "":
                        markdown_lines.append("")
                    
                    if len(doc.tables) > 1:
                        markdown_lines.append(f"### Table {table_num + 1}\n")
                    
                    # Process table rows
                    for row_idx, row in enumerate(table.rows):
                        cells = [cell.text.strip().replace('\n', ' ').replace('|', '\\|') for cell in row.cells]
                        markdown_lines.append('| ' + ' | '.join(cells) + ' |')
                        
                        # Add separator after header row
                        if row_idx == 0:
                            markdown_lines.append('| ' + ' | '.join(['---'] * len(cells)) + ' |')
                    
                    markdown_lines.append("")  # Add blank line after table
            
            if not markdown_lines or all(not line for line in markdown_lines):
                return "*Empty document*"
            
            return "\n".join(markdown_lines).strip()
            
        except Exception as e:
            raise Exception(f"DOCX extraction failed: {str(e)}")
    
    def txt_to_md(self):
        """
        Convert TXT file to markdown format.
        Simply reads the text content and returns it.
        
        Returns:
            Markdown formatted string representation of the TXT file
        """
        try:
            with open(self.file_path, 'r', encoding='utf-8') as txt_file:
                content = txt_file.read()
                
                if not content.strip():
                    return "*Empty text file*"
                
                return content.strip()
                
        except Exception as e:
            raise Exception(f"TXT extraction failed: {str(e)}")
    
    def md_to_md(self):
        """
        Convert MD (Markdown) file to markdown format.
        Simply reads the markdown content and returns it as-is.
        
        Returns:
            Markdown content from the file
        """
        try:
            with open(self.file_path, 'r', encoding='utf-8') as md_file:
                content = md_file.read()
                
                if not content.strip():
                    return "*Empty markdown file*"
                
                return content.strip()
                
        except Exception as e:
            raise Exception(f"MD extraction failed: {str(e)}")
    
    def pptx_to_md(self):
        """
        Convert PPTX file to markdown format.
        Extracts text from slides, including titles, content, and notes.
        
        Returns:
            Markdown formatted string representation of the PPTX
        """
        try:
            from pptx import Presentation
            
            prs = Presentation(self.file_path)
            markdown_sections = []
            
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_content = []
                
                # Add slide header
                slide_content.append(f"## Slide {slide_num}\n")
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        
                        # Check if this is a title shape
                        if hasattr(shape, "shape_type") and shape == slide.shapes.title:
                            slide_content.append(f"### {text}\n")
                        else:
                            # Regular text content
                            slide_content.append(text + "\n")
                    
                    # Extract text from tables
                    if shape.has_table:
                        table = shape.table
                        table_rows = []
                        
                        for row in table.rows:
                            cells = [cell.text.strip().replace('|', '\\|') for cell in row.cells]
                            table_rows.append('| ' + ' | '.join(cells) + ' |')
                        
                        if table_rows:
                            # Add separator after first row (header)
                            if len(table_rows) > 1:
                                table_rows.insert(1, '| ' + ' | '.join(['---'] * len(table.columns)) + ' |')
                            slide_content.append('\n'.join(table_rows) + '\n')
                
                # Extract notes if present
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_content.append(f"\n**Notes:** {notes_text}\n")
                
                # Only add slide if it has content
                if len(slide_content) > 1:  # More than just the slide header
                    markdown_sections.append('\n'.join(slide_content))
            
            if not markdown_sections:
                return "*Empty or unreadable PowerPoint presentation*"
            
            return "\n\n".join(markdown_sections)
            
        except Exception as e:
            raise Exception(f"PPTX extraction failed: {str(e)}")
    
    # Placeholder methods for future file type support
    # def image_to_md(self):
    #     """Convert image file to markdown format (with OCR or description)."""
    #     pass


def extract_file_content(file_instance):
    """
    Extract content from a file based on its type using FileExtractor class.
    
    Args:
        file_instance: File model instance
        
    Returns:
        Extracted content as markdown string
    """
    extractor = FileExtractor(file_instance.file.path, file_instance.filename)
    return extractor.extract()


def process_pending_files():
    """
    Background task to process all files with is_processing=True.
    Extracts content, indexes in ChromaDB (for CSV/XLSX), and marks files as processed.
    """
    print("=" * 50)
    print("STARTING BACKGROUND FILE PROCESSING")
    print("=" * 50)
    
    try:
        # Get all files that need processing
        pending_files = File.objects.filter(is_processing=True)
        print(f"Found {pending_files.count()} files to process")
        
        for file_instance in pending_files:
            print(f"\nProcessing: {file_instance.filename}")
            try:                
                # Extract content (keep for backward compatibility)
                content = extract_file_content(file_instance)
                
                # Update database with extracted content
                file_instance.extracted_content = content
                
                # Index file in ChromaDB for RAG (CSV, XLSX, PDF, DOCX, TXT, MD, and PPTX)
                file_extension = os.path.splitext(file_instance.filename)[1].lower()
                if file_extension in ['.csv', '.xlsx', '.xls', '.pdf', '.docx', '.doc', '.txt', '.md', '.pptx', '.ppt']:
                    try:
                        from core.handlers.knowledge import index_file, delete_file_chunks
                        
                        # Delete old chunks first (in case of re-indexing)
                        # This prevents duplicate ID errors
                        print(f"Removing old chunks for file: {file_instance.filename}")
                        delete_file_chunks(str(file_instance.uuid))
                        
                        # Index the file
                        user_id = file_instance.document.user.id
                        folder_name = file_instance.folder.name if file_instance.folder else None
                        num_chunks = index_file(
                            file_path=file_instance.file.path,
                            file_id=str(file_instance.uuid),
                            filename=file_instance.filename,
                            user_id=user_id,
                            folder_name=folder_name
                        )
                        print(f"✓ Indexed {num_chunks} chunks in ChromaDB")
                    except Exception as index_error:
                        print(f"✗ Error indexing in ChromaDB: {str(index_error)}")
                        # Don't fail the whole process if indexing fails
                
                # Mark as processed
                file_instance.is_processing = False
                file_instance.save()
                
                print(f"✓ Successfully processed: {file_instance.filename}")
                
            except Exception as e:
                print(f"✗ Error processing file {file_instance.filename}: {str(e)}")
                import traceback
                traceback.print_exc()
                # Mark as not processing even if there's an error
                file_instance.is_processing = False
                file_instance.extracted_content = f"Error during processing: {str(e)}"
                file_instance.save()
        
        print(f"\n{'=' * 50}")
        print(f"COMPLETED: Processed {pending_files.count()} files")
        print("=" * 50)
        
    except Exception as e:
        print(f"✗ CRITICAL ERROR in background processing: {str(e)}")
        import traceback
        traceback.print_exc()


def start_background_processing():
    """
    Start background file processing in a separate thread.
    """
    print("🚀 Initiating background file processing thread...")
    try:
        thread = threading.Thread(target=process_pending_files, daemon=True)
        thread.start()
        print("✓ Background file processing thread started successfully")
    except Exception as e:
        print(f"✗ Failed to start background thread: {str(e)}")
        import traceback
        traceback.print_exc()
